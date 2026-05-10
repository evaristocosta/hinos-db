"""Convert PPTX files in `slides_adapt/` into text dumps in `slides_txt/`.

Refactored: smaller functions, CLI, path handling and logging.
"""

from glob import glob
import logging
from argparse import ArgumentParser
from pathlib import Path
from typing import Mapping, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from tqdm import tqdm
import unicodedata

LOGGER = logging.getLogger(__name__)
logging.basicConfig(
    level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s"
)


PUNCT_MAP: Mapping[int, str] = {
    # Dashes / hyphens
    0x2010: "-",  # HYPHEN
    0x2011: "-",  # NON-BREAKING HYPHEN
    0x2012: "-",  # FIGURE DASH
    0x2013: "-",  # EN DASH
    0x2014: "-",  # EM DASH
    0x2015: "-",  # HORIZONTAL BAR
    # Quotes
    0x2018: "'",  # LEFT SINGLE QUOTATION MARK
    0x2019: "'",  # RIGHT SINGLE QUOTATION MARK
    0x201B: "'",  # SINGLE HIGH-REVERSED-9 QUOTATION MARK
    0x201C: '"',  # LEFT DOUBLE QUOTATION MARK
    0x201D: '"',  # RIGHT DOUBLE QUOTATION MARK
    0x201F: '"',  # DOUBLE HIGH-REVERSED-9 QUOTATION MARK
    # Misc punctuation / spacing
    0x00A0: " ",  # NO-BREAK SPACE
    0x2007: " ",  # FIGURE SPACE
    0x2009: " ",  # THIN SPACE
    0x2026: "...",  # HORIZONTAL ELLIPSIS
}


def normalize_text(s: str, ascii_only: bool = False) -> str:
    s = unicodedata.normalize("NFKC", s)
    s = "".join(PUNCT_MAP.get(ord(ch), ch) for ch in s)
    if ascii_only:
        s = s.encode("ascii", "ignore").decode("ascii")
    return s


def find_pptx_files(src: str = "slides_adapt") -> List[str]:
    paths = [
        p for p in glob(str(Path(src) / "*.pptx")) if not Path(p).name.startswith("~$")
    ]
    return paths


def extract_pptx_to_txt(
    file_path: str, out_dir: str = "slides_txt", ascii_only: bool = False
) -> Optional[Path]:
    LOGGER.info("Processing file: %s", file_path)
    prs = Presentation(file_path)
    out_path = Path(out_dir)
    out_path.mkdir(parents=True, exist_ok=True)
    dest = out_path / (Path(file_path).stem + ".txt")

    with dest.open("w", encoding="utf-8") as fh:
        for i, slide in tqdm(
            enumerate(prs.slides), desc="Processing slides", total=len(prs.slides)
        ):
            fh.write(f"\nSLIDE_{i}\n")

            for shape in slide.shapes:
                fh.write(f"\nSHAPE_{shape.shape_type}")

                if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    fh.write(f"_AUTO_SHAPE_{shape.auto_shape_type}\n")
                elif shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
                    fh.write("_SHAPE_FREEFORM_")

                if hasattr(shape, "text") and shape.text.strip() != "":
                    fh.write(
                        f"\nSTART_TEXT\n{normalize_text(shape.text, ascii_only)}\nEND_TEXT\n"
                    )

                fh.write(f"\nHEIGHT_{shape.height}\n")
                fh.write(f"TOP_{shape.top}\n")
                fh.write("\n")

            for shape in slide.shapes:
                if hasattr(shape, "text") and normalize_text(
                    shape.text, ascii_only
                ).lower() in ["indice", "índice"]:
                    fh.write("\n__END__\n")

    LOGGER.info("Wrote: %s", dest)
    return dest


def process_all(
    src: str = "slides_adapt", out_dir: str = "slides_txt", ascii_only: bool = False
) -> List[Path]:
    LOGGER.info("Starting pptx2txt conversion: src=%s out=%s", src, out_dir)
    files = find_pptx_files(src)
    LOGGER.info("Found %d files", len(files))
    written: List[Path] = []
    for f in files:
        p = extract_pptx_to_txt(f, out_dir=out_dir, ascii_only=ascii_only)
        if p:
            written.append(p)
    return written


def main(argv: Optional[List[str]] = None) -> int:
    parser = ArgumentParser(description="Convert PPTX slides to text dumps")
    parser.add_argument(
        "--src", default="slides_adapt", help="Source folder with PPTX files"
    )
    parser.add_argument(
        "--out-dir", default="slides_txt", help="Output directory for TXT files"
    )
    parser.add_argument(
        "--ascii-only", action="store_true", help="Force ASCII-only output"
    )
    parser.add_argument("--quiet", action="store_true")
    args = parser.parse_args(list(argv) if argv is not None else None)

    if args.quiet:
        LOGGER.setLevel(logging.WARNING)

    process_all(src=args.src, out_dir=args.out_dir, ascii_only=args.ascii_only)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
